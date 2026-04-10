# Copyright (C) 2026 Intel Corporation
# SPDX-License-Identifier: Apache-2.0

import logging
import shutil
import time
import traceback
from typing import Dict, List, Callable, Optional
from pptx import Presentation
from config import (
    LLAMA_CONFIG,
    TRANSLATION_CONFIG,
    FILE_CONFIG,
    MODEL_PRESETS,
    FONT_SIZE_ADJUSTMENT,
)
from textgen_client import TextGenClient
from formatting_utils import RunFormatting, extract_run_formatting, apply_run_formatting
from text_utils import get_text_width

logger = logging.getLogger(__name__)


class LlamaPPTTranslator:
    """Translate PowerPoint decks while preserving formatting."""

    def __init__(
        self,
        base_url: str | None = None,
        model_preset: str = "qwen",
        target_language: str | None = None,
        source_language: str = "English",
        llama_config: dict | None = None,
        translation_config: dict | None = None,
        file_config: dict | None = None,
        model_presets: dict | None = None,
        font_adjustment: dict | None = None,
    ):
        self.llama_config = llama_config or LLAMA_CONFIG
        self.translation_config = translation_config or TRANSLATION_CONFIG
        self.file_config = file_config or FILE_CONFIG
        self.model_presets = model_presets or MODEL_PRESETS
        self.font_adjustment_settings = font_adjustment or FONT_SIZE_ADJUSTMENT

        self.target_language = target_language if target_language is not None else self.translation_config["target_language"]
        self.source_language = source_language if source_language is not None else self.translation_config.get("source_language", "English")
        self.adjustment_method = self.font_adjustment_settings.get("method", "dynamic")

        merged_model = {**self.llama_config, **self.model_presets.get(model_preset, {})}
        self.base_url = base_url or merged_model["base_url"]
        self.model_config = merged_model

        self.client = TextGenClient(
            base_url=self.base_url,
            model=self.model_config["model"],
            timeout=self.model_config.get("timeout_seconds", 120)
        )

        if not self.client.test_connection():
            raise ConnectionError(f"Cannot connect to text generation service at {self.base_url}")

    # --- Formatting helpers -------------------------------------------------
    def get_fixed_font_multiplier(self) -> float:
        fixed_rules = self.font_adjustment_settings["fixed"]
        key = f"{self.source_language}->{self.target_language}"
        return fixed_rules.get(key, fixed_rules.get("default", 1.0))

    def calculate_paragraph_font_adjustment(self, paragraph_structure: Dict, translated_texts: List[str]) -> float:
        if not translated_texts:
            return 1.0

        original_full = ""
        translated_full = ""
        font_name = "Arial"
        avg_font_size = 0.0
        font_count = 0
        text_index = 0

        for run_data in paragraph_structure["runs"]:
            original_text = run_data["text"]
            if original_text.strip():
                original_full += original_text
                if text_index < len(translated_texts):
                    translated_full += translated_texts[text_index]
                    text_index += 1

                formatting: RunFormatting = run_data["formatting"]
                if formatting.font_size:
                    avg_font_size += formatting.font_size
                    font_count += 1
                if formatting.font_name and font_name == "Arial":
                    font_name = formatting.font_name

        if not original_full.strip() or not translated_full.strip():
            return 1.0

        if font_count > 0:
            avg_font_size = avg_font_size / font_count
        else:
            avg_font_size = 12.0

        original_width = get_text_width(original_full, font_name, avg_font_size)
        translated_width = get_text_width(translated_full, font_name, avg_font_size)
        if original_width == 0:
            return 1.0

        width_ratio = translated_width / original_width
        settings = self.font_adjustment_settings["dynamic"]
        target_ratio = settings["target_fill_ratio"]
        adjustment = target_ratio / width_ratio if width_ratio > 1.0 else 1.0
        adjustment = max(adjustment, settings["min_adjustment"])
        adjustment = min(adjustment, settings["max_adjustment"])

        if abs(adjustment - 1.0) > 0.05:
            logger.debug(
                "Paragraph font adjustment: %.2fx (width ratio: %.2f, orig: %.0fpx, trans: %.0fpx)",
                adjustment,
                width_ratio,
                original_width,
                translated_width,
            )
        return adjustment

    # --- Translation helpers ------------------------------------------------
    def translate_text_batch(self, texts: List[str]) -> List[str]:
        if not texts or all(not text.strip() for text in texts):
            return texts

        cfg = self.translation_config
        text_mapping = []
        filtered_texts: List[str] = []
        for i, text in enumerate(texts):
            if text.strip():
                text_mapping.append((i, len(filtered_texts)))
                filtered_texts.append(text)
            else:
                text_mapping.append((i, -1))

        if not filtered_texts:
            return texts

        numbered_texts = [f"{i}: {text}" for i, text in enumerate(filtered_texts)]

        preservation_rules = ""
        if cfg.get("preserve_proper_nouns", True):
            custom_rules = cfg.get("custom_preservation_rules", [])
            if custom_rules:
                preservation_rules = "\n".join([f"- {rule}" for rule in custom_rules])
            else:
                preservation_rules = (
                    "- Do NOT translate ANY human names (Western, Asian, or any origin - keep exactly as written)\n"
                    "- Do NOT translate company names, brand names, or product names\n"
                    "- Do NOT translate names that appear to be in Chinese/Japanese/Korean characters - they are likely person names\n"
                    "- Keep ALL proper nouns in their original form and original script\n"
                    "- Examples: Keep \"Zhang Wei\", \"李明\", \"Tanaka\", \"John Smith\" all unchanged"
                )

        prompt = (
            f"/no_think\n"
            f"Translate each text to {self.target_language}. \n\n"
            "CRITICAL RULES - READ CAREFULLY:\n"
            f"{preservation_rules}\n"
            "- Output only the translation for each line, no explanations\n\n"
            "Texts to translate:\n"
            f"{chr(10).join(numbered_texts)}\n\n"
            "Output format (number: translation):\n"
            "0: translated text here\n"
            "1: translated text here\n"
            "Do not include brackets or quotes around the translation."
        )

        for attempt in range(cfg["retry_attempts"]):
            try:
                logger.debug("Translation attempt %s", attempt + 1)
                logger.debug("Target language: %s", self.target_language)

                system_prompt = (
                    f"You are a professional translator. Translate ONLY the text to {self.target_language}.\n\n"
                    "CRITICAL INSTRUCTIONS:\n"
                    "- Output ONLY the translated text in the exact format requested\n"
                    "- Do NOT add any explanations, comments, or extra text\n"
                    "- Do NOT add phrases like \"这里是翻译文本\" or \"Here is the translation\"\n"
                    "- Do NOT add any introductory or concluding remarks\n"
                    "- Follow the exact output format: \"index: translation\" for each line"
                )

                context = cfg.get("presentation_context", "")
                if context:
                    system_prompt += (
                        f"\n\nPresentation context: {context}\nUse this context to provide accurate, domain-specific translations."
                    )

                if cfg.get("preserve_proper_nouns", True):
                    system_prompt += (
                        " CRITICAL: NEVER translate human names regardless of their origin (Western, Chinese, Japanese, Korean, etc.). "
                        "If you see names in Chinese characters like 李明, 张伟, keep them exactly as they are. "
                        "If you see Western names like John Smith, keep them unchanged. "
                        "NEVER translate company names or brand names. Keep ALL proper nouns in their original form and script."
                    )

                response_content = self.client.chat_completion(system_prompt, prompt)
                translations = self.parse_translation_response(response_content, filtered_texts)

                result = [""] * len(texts)
                for orig_idx, filtered_idx in text_mapping:
                    if filtered_idx == -1:
                        result[orig_idx] = texts[orig_idx]
                    else:
                        result[orig_idx] = (
                            translations[filtered_idx] if filtered_idx < len(translations) else texts[orig_idx]
                        )
                return result
            except Exception as exc:
                logger.warning("Translation attempt %s failed: %s", attempt + 1, exc)
                if attempt == cfg["retry_attempts"] - 1:
                    logger.error("All translation attempts failed, returning original texts")
                    return texts
                time.sleep(2**attempt)

        return texts

    def parse_translation_response(self, response_content: str, original_texts: List[str]) -> List[str]:
        translations = [""] * len(original_texts)
        commentary_phrases = [
            "这里是翻译文本",
            "这里是译文",
            "以下是翻译",
            "翻译如下",
            "Here is the translation",
            "Translation:",
            "Translated text:",
            "以下是翻译内容",
            "翻译内容如下",
        ]

        logger.debug("Raw LLM response:\n%s", response_content)

        for line in response_content.split("\n"):
            line = line.strip()
            if not line or any(phrase in line for phrase in commentary_phrases):
                continue
            if ":" in line:
                try:
                    idx_str, translation = line.split(":", 1)
                    idx = int(idx_str.strip())
                    if 0 <= idx < len(translations):
                        translation = translation.strip()
                        if translation.startswith("[") and translation.endswith("]"):
                            translation = translation[1:-1].strip()
                        for phrase in commentary_phrases:
                            translation = translation.replace(phrase, "").strip()
                        translations[idx] = translation
                except (ValueError, IndexError):
                    continue

        for i, (original, translated) in enumerate(zip(original_texts, translations)):
            if not translated.strip():
                translations[i] = original

        logger.debug("Parsed translations: %s", translations)
        return translations

    # --- Paragraph rebuild ---------------------------------------------------
    def extract_paragraph_structure(self, paragraph) -> Dict:
        structure = {
            "alignment": getattr(paragraph, "alignment", None),
            "space_before": getattr(paragraph, "space_before", None),
            "space_after": getattr(paragraph, "space_after", None),
            "line_spacing": getattr(paragraph, "line_spacing", None),
            "level": getattr(paragraph, "level", None),
            "runs": [],
        }

        for run in paragraph.runs:
            run_data = {"text": run.text, "formatting": extract_run_formatting(run)}
            structure["runs"].append(run_data)

        return structure

    def rebuild_paragraph_with_formatting(self, paragraph, structure: Dict, translated_texts: List[str]):
        try:
            font_adjustment = 1.0
            if self.translation_config.get("auto_adjust_font_size", False) and self.adjustment_method == "dynamic":
                font_adjustment = self.calculate_paragraph_font_adjustment(structure, translated_texts)
            elif self.adjustment_method == "fixed":
                font_adjustment = self.get_fixed_font_multiplier()

            for _ in range(len(paragraph.runs)):
                if len(paragraph.runs) > 0:
                    paragraph._p.remove(paragraph.runs[0]._r)

            if structure["alignment"] is not None:
                paragraph.alignment = structure["alignment"]
            if structure["space_before"] is not None:
                paragraph.space_before = structure["space_before"]
            if structure["space_after"] is not None:
                paragraph.space_after = structure["space_after"]
            if structure["line_spacing"] is not None:
                paragraph.line_spacing = structure["line_spacing"]
            if structure["level"] is not None:
                paragraph.level = structure["level"]

            text_index = 0
            dynamic_settings = self.font_adjustment_settings["dynamic"]
            for run_data in structure["runs"]:
                original_text = run_data["text"]
                if original_text.strip():
                    new_text = translated_texts[text_index] if text_index < len(translated_texts) else original_text
                    text_index += 1
                else:
                    new_text = original_text

                new_run = paragraph.add_run()
                new_run.text = new_text
                apply_run_formatting(new_run, run_data["formatting"], font_adjustment, dynamic_settings)
        except Exception as exc:  # pragma: no cover - defensive
            logger.error(f"Error rebuilding paragraph: {exc}")
            logger.error(traceback.format_exc())

    # --- Shape processing ----------------------------------------------------
    def process_text_frame(self, text_frame):
        try:
            all_texts: List[str] = []
            paragraph_structures: List[Dict] = []

            for paragraph in text_frame.paragraphs:
                structure = self.extract_paragraph_structure(paragraph)
                paragraph_structures.append(structure)

                paragraph_texts = []
                for run_data in structure["runs"]:
                    if run_data["text"].strip():
                        paragraph_texts.append(run_data["text"])

                all_texts.extend(paragraph_texts)

            if all_texts:
                batch_size = self.translation_config["batch_size"]
                all_translated: List[str] = []
                for i in range(0, len(all_texts), batch_size):
                    batch = all_texts[i : i + batch_size]
                    translated_batch = self.translate_text_batch(batch)
                    all_translated.extend(translated_batch)
            else:
                all_translated = []

            text_index = 0
            for i, paragraph in enumerate(text_frame.paragraphs):
                structure = paragraph_structures[i]
                paragraph_text_count = sum(1 for run_data in structure["runs"] if run_data["text"].strip())
                paragraph_translations = all_translated[text_index : text_index + paragraph_text_count]
                text_index += paragraph_text_count
                self.rebuild_paragraph_with_formatting(paragraph, structure, paragraph_translations)
        except Exception as exc:  # pragma: no cover - defensive
            logger.error(f"Error processing text frame: {exc}")

    def process_shape(self, shape, shape_path: str = ""):
        try:
            shape_name = getattr(shape, "name", "Unknown")
            shape_type = getattr(shape, "shape_type", "Unknown")
            current_path = f"{shape_path}/{shape_name}" if shape_path else shape_name

            logger.debug("Processing shape: %s (Type: %s)", current_path, shape_type)

            if hasattr(shape, "text_frame"):
                try:
                    text_frame = shape.text_frame
                    if text_frame and text_frame.text.strip():
                        logger.debug("Found text in %s: '%s...'", current_path, text_frame.text[:50])
                        self.process_text_frame(text_frame)
                except Exception as exc:
                    logger.debug("Error checking text_frame in %s: %s", current_path, exc)

            if hasattr(shape, "has_table"):
                try:
                    if shape.has_table:
                        logger.debug("Found table in %s", current_path)
                        self.process_table(shape.table)
                except Exception as exc:
                    logger.debug("Error processing table in %s: %s", current_path, exc)

            if hasattr(shape, "shapes"):
                try:
                    logger.debug("Found %s grouped shapes in %s", len(shape.shapes), current_path)
                    for sub_shape in shape.shapes:
                        self.process_shape(sub_shape, current_path)
                except Exception as exc:
                    logger.debug("Error processing grouped shapes in %s: %s", current_path, exc)

            if hasattr(shape, "has_chart"):
                try:
                    if shape.has_chart:
                        logger.debug("Found chart in %s", current_path)
                        chart = shape.chart
                        if hasattr(chart, "title") and chart.title:
                            if hasattr(chart.title, "text_frame") and chart.title.text_frame:
                                self.process_text_frame(chart.title.text_frame)
                except Exception as exc:
                    logger.debug("Error processing chart in %s: %s", current_path, exc)

        except Exception as exc:  # pragma: no cover - defensive
            logger.error(f"Error processing shape {current_path}: {exc}")
            logger.debug(traceback.format_exc())

    def process_notes_slide(self, slide):
        try:
            if not hasattr(slide, "notes_slide") or not slide.notes_slide:
                return
            notes_slide = slide.notes_slide
            if not hasattr(notes_slide, "notes_text_frame") or not notes_slide.notes_text_frame:
                return
            notes_text_frame = notes_slide.notes_text_frame
            if notes_text_frame.text.strip():
                logger.debug("Found speaker notes: '%s...'", notes_text_frame.text[:100])
                self.process_text_frame(notes_text_frame)
            else:
                logger.debug("No speaker notes text found")
        except Exception as exc:  # pragma: no cover - defensive
            logger.debug(f"Error processing speaker notes: {exc}")

    def process_table(self, table):
        logger.debug("Processing table with %s rows", len(table.rows))
        for row in table.rows:
            for cell in row.cells:
                if hasattr(cell, "text_frame") and cell.text_frame:
                    self.process_text_frame(cell.text_frame)

    # --- Public API ----------------------------------------------------------
    def translate_presentation(self, input_path: str, output_path: str, progress_callback: Optional[Callable[[int, int, str], None]] = None):
        logger.info("Starting translation: %s -> %s", input_path, output_path)
        try:
            if self.file_config.get("backup_original"):
                backup_path = input_path.replace(".pptx", "_backup.pptx")
                shutil.copy2(input_path, backup_path)
                logger.info("Created backup: %s", backup_path)

            prs = Presentation(input_path)
            total_slides = len(prs.slides)
            logger.info("Loaded presentation with %s slides", total_slides)

            for i, slide in enumerate(prs.slides):
                logger.info("Processing slide %s/%s", i + 1, total_slides)
                if progress_callback:
                    progress_callback(i + 1, total_slides, f"Processing slide {i + 1}/{total_slides}")

                total_shapes = len(slide.shapes)
                logger.debug("Slide %s has %s shapes", i + 1, total_shapes)

                for j, shape in enumerate(slide.shapes):
                    logger.debug("Processing shape %s/%s", j + 1, total_shapes)
                    self.process_shape(shape)

                if self.translation_config.get("translate_speaker_notes", True):
                    logger.debug("Processing speaker notes for slide %s", i + 1)
                    self.process_notes_slide(slide)

            prs.save(output_path)
            logger.info("✅ Translation completed successfully: %s", output_path)
            return output_path
        except Exception as exc:
            logger.error(f"❌ Translation failed: {exc}")
            raise